import csv
import json
import os
import shutil
import subprocess
import sys
import time
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader

BASE_DIR = Path(__file__).resolve().parent.parent
BACKEND_DIR = BASE_DIR / "backend"
FRONTEND_DIR = BASE_DIR / "frontend"
FRONTEND_DIST = FRONTEND_DIR / "dist"
EXPORTS_DIR = BASE_DIR / "exports"
FINAL_DIR = BASE_DIR / "FINAL_SUBMISSION"
SCREENSHOTS_DIR = FINAL_DIR / "screenshots"
SAMPLE_OUTPUTS_DIR = FINAL_DIR / "sample_outputs"
SOURCE_CODE_DIR = FINAL_DIR / "source_code"

JOB_TITLE = "Senior Python Backend Engineer"
JOB_DESCRIPTION = (
    "We are hiring a Senior Python Backend Engineer with Django, REST APIs, PostgreSQL, and cloud deployment experience."
)


def run_command(command, cwd=None, env=None, capture_output=False):
    print(f"Running: {' '.join(str(p) for p in command)} in {cwd}")
    result = subprocess.run(
        [str(p) for p in command],
        cwd=cwd,
        env=env,
        capture_output=capture_output,
        text=True,
        check=True,
    )
    return result


def ensure_dirs():
    FINAL_DIR.mkdir(parents=True, exist_ok=True)
    SCREENSHOTS_DIR.mkdir(parents=True, exist_ok=True)
    SAMPLE_OUTPUTS_DIR.mkdir(parents=True, exist_ok=True)
    EXPORTS_DIR.mkdir(parents=True, exist_ok=True)


def seed_data():
    run_command([sys.executable, "scripts/seed_demo.py"], cwd=BACKEND_DIR)


def build_frontend():
    npm_executable = shutil.which("npm.cmd") or shutil.which("npm")
    if not npm_executable:
        raise RuntimeError("npm executable not found on PATH")
    run_command([npm_executable, "run", "build"], cwd=FRONTEND_DIR)


def generate_sample_rankings():
    run_command([sys.executable, "scripts/generate_sample_rankings.py"], cwd=BACKEND_DIR)


def create_sample_explanations():
    source_json = EXPORTS_DIR / "sample_ranked_candidates.json"
    if not source_json.exists():
        raise FileNotFoundError(source_json)
    with source_json.open("r", encoding="utf-8") as fh:
        payload = json.load(fh)
    explanations = {
        "job_title": payload.get("job_title"),
        "job_structure": payload.get("job_structure"),
        "results": [
            {
                "candidate_id": item.get("candidate_id"),
                "full_name": item.get("full_name"),
                "semantic_score": item.get("semantic_score"),
                "skill_score": item.get("skill_score"),
                "experience_score": item.get("experience_score"),
                "recruitability_score": item.get("recruitability_score"),
                "llm_score": item.get("llm_score"),
                "final_score": item.get("final_score"),
                "semantic_explanation": item.get("semantic_explanation"),
                "skill_explanation": item.get("skill_explanation"),
                "experience_explanation": item.get("experience_explanation"),
                "behavioral_explanation": item.get("behavioral_explanation"),
                "gemini_explanation": item.get("gemini_explanation"),
                "strengths": item.get("strengths"),
                "weaknesses": item.get("weaknesses"),
                "missing_skills": item.get("missing_skills"),
                "recommendation": item.get("recruiter_ai", {}).get("recommendation"),
            }
            for item in payload.get("results", [])
        ],
    }
    out_path = SAMPLE_OUTPUTS_DIR / "sample_explanations.json"
    with out_path.open("w", encoding="utf-8") as fh:
        json.dump(explanations, fh, indent=2)
    print(f"Wrote sample explanations to {out_path}")
    return out_path


def copy_sample_outputs():
    ranked_csv = EXPORTS_DIR / "ranked_candidates.csv"
    if ranked_csv.exists():
        shutil.copy(ranked_csv, FINAL_DIR / "ranked_candidates.csv")
    sample_csv = SAMPLE_OUTPUTS_DIR / "sample_ranked_output.csv"
    with sample_csv.open("w", newline="", encoding="utf-8") as fh:
        writer = csv.writer(fh)
        writer.writerow([
            "candidate_id",
            "full_name",
            "semantic_score",
            "skill_score",
            "experience_score",
            "recruitability_score",
            "llm_score",
            "final_score",
            "recommendation",
        ])
        sample_json = EXPORTS_DIR / "sample_ranked_candidates.json"
        with sample_json.open("r", encoding="utf-8") as source_f:
            data = json.load(source_f)
        for item in data.get("results", [])[:5]:
            writer.writerow([
                item.get("candidate_id"),
                item.get("full_name"),
                item.get("semantic_score"),
                item.get("skill_score"),
                item.get("experience_score"),
                item.get("recruitability_score"),
                item.get("llm_score"),
                item.get("final_score"),
                item.get("recruiter_ai", {}).get("recommendation", ""),
            ])
    print(f"Wrote sample ranked output CSV to {sample_csv}")
    return sample_csv


def draw_architecture_diagram():
    diagram_path = FINAL_DIR / "architecture_diagram.png"
    width, height = 1800, 900
    image = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(image)
    font_path = None
    try:
        from PIL import ImageFont
        font_path = ImageFont.load_default()
        font = font_path
    except Exception:
        font = None
    box_fill = "#dbeafe"
    outline = "#1d4ed8"
    text_color = "#0f172a"

    boxes = [
        ((120, 120), (520, 240), "Job Understanding"),
        ((640, 120), (1040, 240), "Semantic Search"),
        ((1160, 120), (1560, 240), "Hybrid Ranking"),
        ((120, 340), (520, 460), "Skill Engine"),
        ((640, 340), (1040, 460), "Experience Engine"),
        ((1160, 340), (1560, 460), "Behavioral Engine"),
        ((640, 580), (1040, 700), "Gemini Recruiter Agent"),
        ((1160, 580), (1560, 700), "Explainability"),
    ]

    for top_left, bottom_right, text in boxes:
        draw.rectangle([top_left, bottom_right], fill=box_fill, outline=outline, width=4)
        text_x = top_left[0] + 20
        text_y = top_left[1] + 20
        draw.text((text_x, text_y), text, fill=text_color, font=font)

    arrows = [
        ((520, 180), (640, 180)),
        ((1040, 180), (1160, 180)),
        ((320, 240), (320, 340)),
        ((860, 240), (860, 340)),
        ((1360, 240), (1360, 340)),
        ((860, 460), (860, 580)),
        ((1360, 460), (1360, 580)),
        ((1040, 640), (1160, 640)),
    ]
    for start, end in arrows:
        draw.line([start, end], fill=outline, width=5)
        draw.polygon([end, (end[0]-15, end[1]-10), (end[0]-15, end[1]+10)], fill=outline)

    image.save(diagram_path)
    print(f"Created architecture diagram at {diagram_path}")
    return diagram_path


def capture_screenshots():
    from playwright.sync_api import sync_playwright

    backend_proc = subprocess.Popen(
        [sys.executable, "manage.py", "runserver", "8000"], cwd=BACKEND_DIR, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
    )
    static_proc = subprocess.Popen(
        [sys.executable, "-m", "http.server", "4174", "--directory", str(FRONTEND_DIST)],
        cwd=BASE_DIR,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    try:
        import urllib.request
        for _ in range(30):
            try:
                urllib.request.urlopen("http://127.0.0.1:8000/api/candidates/", timeout=1)
                urllib.request.urlopen("http://127.0.0.1:4174", timeout=1)
                break
            except Exception:
                time.sleep(1)
        else:
            raise RuntimeError("Servers did not start in time")

        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page()
            page.goto("http://127.0.0.1:4174", wait_until="networkidle")
            page.screenshot(path=SCREENSHOTS_DIR / "dashboard_home.png", full_page=True)
            page.fill("input", JOB_TITLE)
            page.fill("textarea", JOB_DESCRIPTION)
            page.click("text=Analyze & Rank")
            page.wait_for_selector("text=Ranked candidates", timeout=20000)
            page.wait_for_timeout(3000)
            page.screenshot(path=SCREENSHOTS_DIR / "dashboard_ranked.png", full_page=True)
            browser.close()
        print(f"Captured screenshots in {SCREENSHOTS_DIR}")
    finally:
        backend_proc.terminate()
        static_proc.terminate()
        backend_proc.wait()
        static_proc.wait()


def create_presentation(diagram_path):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "TalentMind AI — Candidate Ranking"
    slide.placeholders[1].text = "Hybrid recruiter intelligence with Gemini, FAISS search, explainability, and dashboard export."

    def add_bullet_slide(title, bullets):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.clear()
        for item in bullets:
            p = body.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
        return slide

    slides = [
        ("Problem Statement", [
            "Hiring teams struggle to rank candidates consistently and quickly.",
            "Traditional ATS systems prioritize resumes, not candidate fit signals.",
            "Need an AI-first recruiter assistant that is explainable and fast.",
        ]),
        ("Existing ATS Limitations", [
            "Keyword matching misses candidate intent and skill relevance.",
            "Little to no insight into candidate behaviors or recruiter fit.",
            "Ranking is rarely explainable or semantically aware.",
        ]),
        ("Proposed Solution", [
            "A hybrid ranking platform combining embeddings, skills, experience, and LLM recruiter signals.",
            "End-to-end dashboard with job understanding, candidate explainability, and exports.",
            "Fallback logic keeps the system robust when Gemini is unavailable.",
        ]),
        ("System Architecture", [
            "React dashboard communicates with Django REST APIs.",
            "Backend computes embeddings, ranking, and Gemini recruiter analysis.",
            "Results are exported to CSV and submission-ready formats.",
        ]),
        ("Data Pipeline", [
            "Job description analysis produces structured required skills and behavioral themes.",
            "Candidate profiles are embedded for semantic matching.",
            "Scores are combined through a weighted hybrid formula.",
        ]),
        ("Job Understanding Engine", [
            "Extracts key skills, experience focus, and behavioral expectations.",
            "Supports recruiter input through structured JD analysis.",
            "Feeds downstream ranking and explainability layers.",
        ]),
        ("Semantic Search + FAISS", [
            "Candidate embeddings are indexed for fast nearest-neighbor retrieval.",
            "Semantic score captures intent and overall resume fit.",
            "FAISS enables scalable similarity search in production.",
        ]),
        ("Skill Intelligence Engine", [
            "Matches required skills against candidate skill profiles.",
            "Generates strengths and gap detection.",
            "Contributes to explainable role fit scoring.",
        ]),
        ("Experience Intelligence Engine", [
            "Evaluates years of experience and career trajectory.",
            "Looks for stability, growth, and role alignment.",
            "Provides experience-specific explainability text.",
        ]),
        ("Behavioral Intelligence Engine", [
            "Assesses recruitability signals like response and completion rates.",
            "Captures open-to-work status and profile completeness.",
            "Improves ranking with practical hiring readiness metrics.",
        ]),
        ("Gemini Recruiter Agent", [
            "Invokes Gemini for qualitative recruiter analysis.",
            "Produces an LLM fit score, recommendation, strengths, weaknesses, and summary.",
            "Falls back to deterministic logic when Gemini is unavailable.",
        ]),
        ("Hybrid Ranking Formula", [
            "final_score = 0.35*semantic + 0.25*skill + 0.15*experience + 0.15*recruitability + 0.10*llm.",
            "Combines structured and generative signals for balanced ranking.",
            "Supports explainable decisions at each scoring stage.",
        ]),
        ("Explainable AI", [
            "Each candidate receives semantic, skill, experience, behavioral, and Gemini explanations.",
            "Strengths, weaknesses, and missing skills are surfaced.",
            "Enables recruiter trust and fast shortlist decisions.",
        ]),
    ]
    for title, bullets in slides:
        add_bullet_slide(title, bullets)

    arch_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(arch_slide_layout)
    slide.shapes.title.text = "Architecture Diagram"
    slide.shapes.add_picture(str(diagram_path), Inches(0.5), Inches(1.4), width=Inches(9))

    screenshot_files = sorted(SCREENSHOTS_DIR.glob("*.png"))
    for screenshot in screenshot_files:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = screenshot.stem.replace("_", " ").title()
        slide.shapes.add_picture(str(screenshot), Inches(1), Inches(1.5), width=Inches(8))

    result_slide = add_bullet_slide("Results & Impact", [
        "Generated submission-ready ranked outputs and candidate explainability data.",
        "Automation script creates the final submission package and screenshots.",
        "Demonstrated a robust fallback for Gemini API availability.",
    ])
    add_bullet_slide("Future Enhancements", [
        "Add enterprise-grade persona and diversity-aware ranking.",
        "Extend candidate sourcing with resume parsing and ATS ingestion.",
        "Build a multi-user role-based dashboard and interview insights.",
    ])

    pptx_path = FINAL_DIR / "TalentMind_AI_Presentation.pptx"
    prs.save(str(pptx_path))
    print(f"Saved presentation to {pptx_path}")
    return pptx_path


def create_pdf(pptx_path, diagram_path):
    pdf_path = FINAL_DIR / "TalentMind_AI_Presentation.pdf"
    c = canvas.Canvas(str(pdf_path), pagesize=letter)
    width, height = letter
    margin = 50

    def draw_page(title, lines, image_path=None):
        c.setFont("Helvetica-Bold", 24)
        c.drawString(margin, height - margin, title)
        c.setFont("Helvetica", 12)
        y = height - margin - 40
        for line in lines:
            c.drawString(margin, y, line)
            y -= 18
        if image_path and Path(image_path).exists():
            img = ImageReader(str(image_path))
            c.drawImage(img, margin, 60, width=width - margin * 2, height=height / 2, preserveAspectRatio=True)
        c.showPage()

    pages = [
        ("Problem Statement", [
            "Hiring teams need faster, explainable candidate ranking.",
            "Current ATS solutions do not combine semantic, behavioral, and LLM recruiter intelligence.",
        ], None),
        ("Solution Overview", [
            "TalentMind AI blends embeddings, skill and experience scoring, behavioral signals, and Gemini recruiter insights.",
            "The output is a ranked candidate shortlist with detailed explanations.",
        ], str(diagram_path)),
        ("Dashboard Screenshots", ["Dashboard home and ranked candidate views are captured automatically."], str(SCREENSHOTS_DIR / "dashboard_home.png")),
        ("Dashboard Results", ["The platform generates submission-ready CSV exports and JSON explainability data."], str(SCREENSHOTS_DIR / "dashboard_ranked.png")),
    ]
    for title, lines, image_path in pages:
        draw_page(title, lines, image_path)
    c.save()
    print(f"Saved PDF to {pdf_path}")
    return pdf_path


def copy_source_code():
    if SOURCE_CODE_DIR.exists():
        shutil.rmtree(SOURCE_CODE_DIR)
    SOURCE_CODE_DIR.mkdir(parents=True, exist_ok=True)
    for name in ["backend", "frontend", "docs"]:
        src = BASE_DIR / name
        dst = SOURCE_CODE_DIR / name
        shutil.copytree(src, dst, dirs_exist_ok=True, ignore=shutil.ignore_patterns("__pycache__", "*.pyc", "node_modules", "dist", ".git", ".venv", "env", "db.sqlite3"))
    for file_name in ["requirements.txt", "package.json", "README.md"]:
        src = BASE_DIR / file_name
        if src.exists():
            shutil.copy(src, SOURCE_CODE_DIR / file_name)


def create_final_readme():
    final_readme = FINAL_DIR / "README.md"
    final_readme.write_text(
        "# TalentMind AI Final Submission\n\n"
        "## Installation\n"
        "1. Install backend dependencies: `pip install -r requirements.txt`\n"
        "2. Install frontend dependencies: `npm install`\n\n"
        "## Setup\n"
        "1. Run `python manage.py migrate` in `backend/`.\n"
        "2. Seed sample data: `python scripts/seed_demo.py`\n"
        "3. Set `GEMINI_API_KEY` if available for LLM recruiter signals.\n\n"
        "## Run Commands\n"
        "- Backend: `python manage.py runserver`\n"
        "- Frontend: `npm run dev`\n"
        "- Generate sample artifacts: `python scripts/generate_submission_package.py`\n\n"
        "## API Endpoints\n"
        "- POST /api/analyze-jd/\n"
        "- POST /api/rank/\n"
        "- GET /api/candidates/\n"
        "- GET /api/rankings/\n"
        "- GET /api/rankings/export/?job_id=<id>&format=submission\n\n"
        "## Project Structure\n"
        "- backend/: Django backend and ranking services\n"
        "- frontend/: React + Tailwind dashboard\n"
        "- docs/: architecture and presentation notes\n"
        "- exports/: generated CSV and JSON outputs\n"
        "- FINAL_SUBMISSION/: final packaged submission artifacts\n\n"
        "## Screenshots\n"
        "- screenshots/dashboard_home.png\n"
        "- screenshots/dashboard_ranked.png\n"
    )
    print(f"Created final README at {final_readme}")
    return final_readme


def package_submission():
    zip_path = BASE_DIR / "TalentMind_AI_Final_Submission.zip"
    if zip_path.exists():
        zip_path.unlink()
    shutil.make_archive(str(zip_path.with_suffix("")), "zip", root_dir=FINAL_DIR)
    print(f"Created submission ZIP at {zip_path}")
    return zip_path


def create_checklist():
    checklist = FINAL_DIR / "submission_checklist.md"
    checklist.write_text(
        "# Submission Checklist\n\n"
        "- [x] Source code included under source_code/\n"
        "- [x] TalentMind_AI_Presentation.pptx included\n"
        "- [x] TalentMind_AI_Presentation.pdf included\n"
        "- [x] ranked_candidates.csv included\n"
        "- [x] sample_ranked_output.csv included\n"
        "- [x] sample_explanations.json included\n"
        "- [x] screenshots included\n"
        "- [x] README.md included\n"
        "- [x] ZIP archive created\n"
    )
    print(f"Created checklist at {checklist}")
    return checklist


if __name__ == "__main__":
    ensure_dirs()
    seed_data()
    build_frontend()
    generate_sample_rankings()
    create_sample_explanations()
    copy_sample_outputs()
    diagram_path = draw_architecture_diagram()
    capture_screenshots()
    pptx_path = create_presentation(diagram_path)
    create_pdf(pptx_path, diagram_path)
    copy_source_code()
    create_final_readme()
    package_submission()
    create_checklist()
    print("Submission package generation complete.")
